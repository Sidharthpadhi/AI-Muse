import concurrent
import datetime
import chardet # If used, ensure it's imported
import streamlit as st
import easyocr
import torch
import numpy as np
import pandas as pd
from PIL import Image
from io import BytesIO
import os
import pdfplumber
from pptx import Presentation
from docx import Document
from concurrent.futures import ProcessPoolExecutor, TimeoutError as FutureTimeoutError
import re
import time
import logging
import multiprocessing # Crucial import

# --- Setup Logging ---
# Configure once at the start. Logs to console where Streamlit is run.
# Set level to logging.DEBUG for more verbose output during development.
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(process)d - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# --- START OF YOUR EXISTING USER FUNCTIONS ---
# IMPORTANT: All these functions must be defined at the top level of this script.
# For brevity, I'm not re-listing all of them, but ensure they are here.
# Example:
torch.classes.__path__ = []
@st.cache_resource
def load_easyocr_model():
    logger.info(f"PID {os.getpid()}: Loading EasyOCR model for this process.")
    # Consider adding model_storage_directory and user_network_directory if needed
    # for specific EasyOCR deployment scenarios, or if it helps with pickling/loading.
    reader = easyocr.Reader(['en', 'hi'], gpu=False, verbose=False) # gpu=False for CPU, verbose=False for cleaner logs
    logger.info(f"PID {os.getpid()}: EasyOCR model loaded.")
    return reader

def extract_text_easyocr(pil_img):
    reader = load_easyocr_model() # This will get the cached model for the current process
    cv_img = np.array(pil_img.convert("RGB"))
    results = reader.readtext(cv_img, detail=0) # detail=0 for just text
    return "\n".join(results)

def extract_data_from_pdf(pdf_path_or_buffer):
    extracted_text = ""
    try:
        with pdfplumber.open(pdf_path_or_buffer) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    extracted_text += text + "\n"
    except Exception as e:
        logger.error(f"Error extracting PDF content: {e}", exc_info=False) # exc_info=False for cleaner UI error
        return f"Error extracting PDF: {str(e)}" # Return simple error string
    return extracted_text

def extract_text_from_pptx(file_path_or_buffer):
    try:
        prs = Presentation(file_path_or_buffer)
        text_runs = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
        return '\n'.join(text_runs)
    except Exception as e:
        logger.error(f"Error extracting PPTX content: {e}", exc_info=False)
        return f"Error extracting PPTX: {str(e)}"

def extract_text_from_docx(file_path_or_buffer):
    try:
        doc = Document(file_path_or_buffer)
        full_text = [para.text for para in doc.paragraphs]
        return '\n'.join(full_text)
    except Exception as e:
        logger.error(f"Error extracting DOCX content: {e}", exc_info=False)
        return f"Error extracting DOCX: {str(e)}"

def extract_text_plain(file_like_object):
    encodings = ['utf-8', 'latin-1', 'iso-8859-1', 'cp1252']
    file_like_object.seek(0)
    file_content = file_like_object.read()
    for encoding in encodings:
        try:
            return file_content.decode(encoding)
        except UnicodeDecodeError:
            continue
    logger.warning("extract_text_plain: Unable to decode with available encodings.")
    return "Error: Unable to decode text file with tried encodings."

def preprocess_text(text):
    if text is not None and isinstance(text, str):
        text = text.strip()
        text = re.sub(r'\s+', ' ', text)
        text = ''.join([c for c in text if c.isprintable()])
        # More permissive regex for allowed characters, adjust as needed
        allowed_chars = r"A-Za-z0-9\s.,:;!?@#$%^&*()[]{}_+=\-/|\\'\"<>`~"
        text = re.sub(r"[^" + re.escape(allowed_chars) + r"]", "", text) # Use re.escape for safety
    elif text is None:
        return ""
    return text

def extract_pan_details(text):
    if not text or not isinstance(text, str): return {"Error": "Invalid input text for PAN"}
    details = {}; text_upper = text.upper()
    if re.search(r"INCOME TAX DEPARTMENT", text_upper): details["Department"] = "INCOME TAX DEPARTMENT"
    if re.search(r"GOVT\.?\s+OF\s+INDIA", text_upper): details["Issuer"] = "GOVT. OF INDIA"
    if re.search(r"PERMANENT ACCOUNT NUMBER CARD", text_upper): details["Document_Title"] = "PERMANENT ACCOUNT NUMBER CARD"
    pan_match = re.search(r"\b[A-Z]{5}[0-9]{4}[A-Z]\b", text_upper);
    if pan_match: details["PAN_Number"] = pan_match.group(0)
    # Simplified name extraction - look for common patterns near "Name"
    name_patterns = [
        r"(?:APPLICANT\s+)?NAME\s*[:\-]?\s*([A-Z\s.]+)", # Allows dots in names
        r"\b([A-Z][A-Z\s.]+)\nFATHER", # Name often appears before Father's Name
    ]
    for pattern in name_patterns:
        name_match = re.search(pattern, text_upper)
        if name_match and name_match.group(1).strip():
            details["Name"] = name_match.group(1).strip(); break
    father_match = re.search(r"FATHER(?:'S)?\s+NAME\s*[:\-]?\s*([A-Z\s.]+)", text_upper)
    if father_match: details["Father_Name"] = father_match.group(1).strip()
    dob_match = re.search(r"\b(\d{2}/\d{2}/\d{4})\b", text_upper) # Capture group for clarity
    if dob_match: details["Date_of_Birth"] = dob_match.group(1)
    details["Signature_Mentioned"] = "Yes" if re.search(r"SIGNATURE", text_upper) else "No"
    return details

def extract_aadhaar_details(text):
    if not text or not isinstance(text, str): return {"Error": "Invalid input text for Aadhaar"}
    details = {}; text_upper = text.upper()
    if "UNIQUE IDENTIFICATION AUTHORITY OF INDIA" in text_upper: details["Issuer"] = "UIDAI"
    if "AADHAAR" in text_upper: details["Document_Title"] = "AADHAAR CARD"
    aadhaar_match = re.search(r"\b(\d{4}\s?\d{4}\s?\d{4})\b", text_upper)
    if aadhaar_match: details["Aadhaar_Number"] = aadhaar_match.group(1)
    enrol_match = re.search(r"ENROLMENT\s*NO\.?\s*[:\-]?\s*(\d{4}/\d{5}/\d{5})", text_upper, re.IGNORECASE)
    if enrol_match: details["Enrolment_No"] = enrol_match.group(1)
    name_match = re.search(r"NAME\s*[:\-]?\s*([A-Z\s.]+)", text_upper)
    if name_match: details["Name"] = name_match.group(1).strip()
    dob_match = re.search(r"DATE\s+OF\s+BIRTH\s*[:\-]?\s*(\d{2}/\d{2}/\d{4})", text_upper, re.IGNORECASE)
    if not dob_match: dob_match = re.search(r"DOB\s*[:\-]?\s*(\d{2}/\d{2}/\d{4})", text_upper, re.IGNORECASE)
    if dob_match: details["Date_of_Birth"] = dob_match.group(1)
    gender_match = re.search(r"\b(MALE|FEMALE|TRANSGENDER)\b", text_upper)
    if gender_match: details["Gender"] = gender_match.group(1)
    # More robust address extraction might be needed depending on format variations
    address_keywords = ["ADDRESS:", "ADDRESS :", "पता:", "पता :"]
    address_lines = []
    found_address_keyword = False
    for line in text_upper.splitlines():
        if any(kw in line for kw in address_keywords):
            found_address_keyword = True
            line_after_keyword = line.split(maxsplit=1)[-1] # Get text after keyword
            if line_after_keyword and not any(kw.startswith(line_after_keyword.strip()) for kw in address_keywords): # Ensure it's not just the keyword again
                address_lines.append(line_after_keyword.strip())
            continue
        if found_address_keyword and line.strip() and not re.match(r"^\d{6}$", line.strip()): # Add line if not just pincode
            address_lines.append(line.strip())
            if len(address_lines) >= 3: break # Limit address lines
        elif found_address_keyword and not line.strip(): # Empty line might end address block
            break
    if address_lines: details["Address"] = " ".join(address_lines).title()
    return details

# ... (Your other extract_..._details functions, identify_document_type, classify_document, etc. ensure they are robust)
# For brevity, I'll assume they are defined here and are top-level.
# Make sure to test them individually for robustness.
# Example for Voter ID (ensure it's robust)
def extract_voter_details(text):
    if not text or not isinstance(text, str): return {"Error": "Invalid input text for Voter ID"}
    details = {}; text_upper = text.upper()
    epic_match = re.search(r"\b([A-Z]{3}[0-9]{7})\b", text_upper) # Common format
    if not epic_match: epic_match = re.search(r"\b([A-Z]{2}/\d{2}/\d{3}/\d{6})\b", text_upper) # Another format
    if epic_match: details["Voter_ID"] = epic_match.group(1)

    if "ELECTION COMMISSION OF INDIA" in text_upper: details["Issuer"] = "ECI"
    if "ELECTOR PHOTO IDENTITY CARD" in text_upper: details["Document_Title"] = "ELECTOR PHOTO IDENTITY CARD"
    
    name_match = re.search(r"(?:NAME|ELECTOR'?S NAME)\s*[:\-]?\s*([A-Z\s.]+)", text_upper, re.IGNORECASE)
    if name_match: details["Name"] = name_match.group(1).strip()
    
    rel_types = ["FATHER", "HUSBAND", "MOTHER"]
    for rel_type in rel_types:
        relative_match = re.search(rf"(?:{rel_type}'?S NAME)\s*[:\-]?\s*([A-Z\s.]+)", text_upper, re.IGNORECASE)
        if relative_match:
            details[f"{rel_type}_Name"] = relative_match.group(1).strip(); break
            
    sex_match = re.search(r"SEX\s*[:\-]?\s*(MALE|FEMALE|OTHER|THIRD GENDER)", text_upper, re.IGNORECASE)
    if sex_match: details["Sex"] = sex_match.group(1)
    
    age_match = re.search(r"AGE\s*(?:AS ON .+)?\s*[:\-]?\s*(\d+)\s*YEARS", text_upper, re.IGNORECASE)
    if age_match: details["Age"] = age_match.group(1)
    
    # विधानसभा क्षेत्र / Assembly Constituency
    ac_match = re.search(r"(\d+)\s*-\s*([A-Za-z\s&(),.-]+)\s*(?:AC)?", text_upper)
    if ac_match: details["Assembly_Constituency"] = f"{ac_match.group(1)}-{ac_match.group(2).strip()}"
    return details

AADHAAR_REGEX = re.compile(r'\b\d{4}\s?\d{4}\s?\d{4}\b', re.IGNORECASE)
PAN_REGEX = re.compile(r'\b[A-Z]{5}[0-9]{4}[A-Z]\b', re.IGNORECASE)
VOTER_ID_REGEX = re.compile(r'\b[A-Z]{3}[0-9]{7}\b', re.IGNORECASE) # One common format
PASSPORT_REGEX = re.compile(r'\b[A-Z][0-9]{7}\b', re.IGNORECASE)
DL_REGEX_STRICT = re.compile(r'\b(?:[A-Z]{2}[-\s]?\d{2}[-\s]?\d{4}[-\s]?\d{7}|[A-Z]{2}[-\s]?\d{13})\b', re.IGNORECASE)
DOCUMENT_PATTERNS = {
    "Aadhaar": (AADHAAR_REGEX, {"AADHAAR", "UNIQUE IDENTIFICATION AUTHORITY OF INDIA", "GOVERNMENT OF INDIA", "UIDAI"}),
    "PAN": (PAN_REGEX, {"INCOME TAX DEPARTMENT", "PERMANENT ACCOUNT NUMBER CARD", "PAN CARD", "GOVT. OF INDIA"}),
    "Voter ID": (VOTER_ID_REGEX, {"VOTER ID", "EPIC", "ELECTION COMMISSION OF INDIA", "ELECTOR PHOTO IDENTITY CARD"}),
    "Passport": (PASSPORT_REGEX, {"PASSPORT", "REPUBLIC OF INDIA", "MINISTRY OF EXTERNAL AFFAIRS"}),
    "Driving License": (DL_REGEX_STRICT, {"DRIVING LICENCE", "TRANSPORT DEPARTMENT", "LICENSE NUMBER", "DL NO", "UNION OF INDIA"}),
    "Medical Document": (None, {"HOSPITAL", "PATIENT", "DOCTOR", "CLINIC", "MEDICAL", "DIAGNOSIS", "PRESCRIPTION", "RX", "REPORT", "LABORATORY"}),
    "Financial Document": (None, {"BANK", "ACCOUNT", "STATEMENT", "INVOICE", "BILL", "FINANCIAL", "TRANSACTION", "RECEIPT", "TAX"})
}
def identify_document_type(text_upper):
    if not text_upper: return "Unknown"
    def _keyword_match_count(keywords, text_content): return sum(1 for keyword in keywords if keyword in text_content)
    detected_scores = {}
    for doc_type, (regex, keywords) in DOCUMENT_PATTERNS.items():
        score = 0
        if regex and regex.search(text_upper): score += 5 # Strong weight for regex match
        score += _keyword_match_count(keywords, text_upper) # Add 1 for each keyword found
        if score > 0: detected_scores[doc_type] = score
    if not detected_scores: return "Unknown"
    # Prioritize types with strong regex and at least one keyword, or high keyword count
    best_match = "Unknown"
    max_score = 0
    for doc_type, score in detected_scores.items():
        regex, keywords = DOCUMENT_PATTERNS[doc_type]
        # Give preference if regex matches AND there's at least one keyword
        if regex and regex.search(text_upper) and _keyword_match_count(keywords, text_upper) > 0:
            # Boost score for regex + keyword combination
            current_effective_score = score + 2 # Arbitrary boost
        else:
            current_effective_score = score
        
        if current_effective_score > max_score:
            max_score = current_effective_score
            best_match = doc_type
        # Tie-breaking: if scores are equal, prefer one with regex match
        elif current_effective_score == max_score and regex and regex.search(text_upper):
            if not (DOCUMENT_PATTERNS[best_match][0] and DOCUMENT_PATTERNS[best_match][0].search(text_upper)):
                best_match = doc_type

    # Minimum threshold for considering a document identified
    if max_score < 2 and best_match != "Unknown": # e.g., only one weak keyword found
         # Check if it's a very generic term like "REPORT" for Medical or "RECEIPT" for Financial
        if best_match in ["Medical Document", "Financial Document"] and max_score < 3: # Higher threshold for these
            return "Unknown"
    
    return best_match

def _keyword_match_count_global(keywords, text): return sum(1 for keyword in keywords if keyword in text) # Ensure this helper is available globally

def classify_document(extracted_details_dict, text_upper): # Takes extracted fields dict
    num_extracted_fields = len(extracted_details_dict) if isinstance(extracted_details_dict, dict) and "Error" not in extracted_details_dict else 0

    # Use identified type if available, otherwise try to infer
    identified_type = extracted_details_dict.get("Document_Title_Identified", identify_document_type(text_upper))

    if identified_type == "Aadhaar" and num_extracted_fields >= 2: return "Aadhaar", "Government ID", "High"
    if identified_type == "PAN" and num_extracted_fields >= 2: return "PAN", "Government ID", "High"
    if identified_type == "Voter ID" and num_extracted_fields >= 1: return "Voter ID", "Government ID", "High" # Voter ID number itself is enough
    if identified_type == "Passport" and num_extracted_fields >= 2: return "Passport", "Government ID", "High"
    if identified_type == "Driving License" and num_extracted_fields >= 2: return "Driving License", "Government ID", "High"
    
    # For keyword-based types, rely more on keyword counts in text_upper
    medical_keywords = DOCUMENT_PATTERNS["Medical Document"][1]
    if _keyword_match_count_global(medical_keywords, text_upper) >= 2: # Needs a couple of medical terms
        return "Medical Document", "Medical Record", "Medium" # Changed type string for clarity

    financial_keywords = DOCUMENT_PATTERNS["Financial Document"][1]
    if _keyword_match_count_global(financial_keywords, text_upper) >= 2: # Needs a couple of financial terms
        return "Financial Document", "Financial Record", "Medium" # Changed type string

    return "Unknown", "Unclassified", "Low" # Default

def mask_data(details_dict):
    masked_details = {}
    if not isinstance(details_dict, dict): return {}
    sensitive_keys = ["PAN_Number", "Aadhaar_Number", "Voter_ID", "Passport_Number", "License_Number", "Account_Number"]
    for key, value in details_dict.items():
        if value and isinstance(value, str):
            if key in sensitive_keys and len(value) > 4: # Mask sensitive numbers more thoroughly
                masked_details[key] = value[:2] + 'X' * (len(value) - 4) + value[-2:]
            elif len(value) > 2 and key != "Document_Title": # General masking for other string values
                masked_details[key] = value[0] + 'X' * (len(value) - 2) + value[-1]
            else:
                masked_details[key] = value # Short strings or titles unmasked
        else:
            masked_details[key] = value
    return masked_details

classification_options = { # This is for the multiselect, not direct classification rules
    "Aadhaar": "Government ID", "PAN": "Government ID", "Voter ID": "Government ID",
    "Passport": "Government ID", "Driving License": "Government ID",
    "Medical Document": "Medical Document", "Financial Document": "Financial Document",
}
def process_excel_csv(file_like_object):
    try:
        # Try reading as Excel first, as it's more specific
        df = pd.read_excel(file_like_object)
        logger.info("Successfully processed as Excel.")
    except Exception as excel_error:
        logger.warning(f"Failed to read as Excel ({excel_error}), trying CSV.")
        try:
            file_like_object.seek(0) # Reset buffer position for CSV reading
            # Attempt to decode as UTF-8, then fall back to chardet or latin-1 for CSV
            try:
                content = file_like_object.read().decode('utf-8')
                df = pd.read_csv(BytesIO(content.encode('utf-8')))
            except UnicodeDecodeError:
                file_like_object.seek(0)
                content_bytes = file_like_object.read()
                detected_encoding = chardet.detect(content_bytes)['encoding']
                if detected_encoding:
                    logger.info(f"Chardet detected encoding: {detected_encoding} for CSV.")
                    df = pd.read_csv(BytesIO(content_bytes), encoding=detected_encoding)
                else: # Fallback to latin-1
                    logger.warning("Chardet couldn't detect encoding, trying latin-1 for CSV.")
                    df = pd.read_csv(BytesIO(content_bytes), encoding='latin-1')
            logger.info("Successfully processed as CSV.")
        except Exception as csv_error:
            logger.error(f"Failed to process as Excel or CSV. Excel error: {excel_error}, CSV error: {csv_error}")
            raise ValueError("Error reading Excel/CSV: File format not recognized or corrupted.") from csv_error
    return df

def save_to_excel(report_data_list_of_dicts):
    if not report_data_list_of_dicts:
        logger.info("No data provided to save_to_excel. Creating empty report.")
        df = pd.DataFrame([{"Message": "No data processed or available for report."}])
    else:
        df = pd.DataFrame(report_data_list_of_dicts)
    output = BytesIO()
    try:
        with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
            df.to_excel(writer, index=False, sheet_name="ScanReport")
        logger.info("Excel report generated successfully.")
        return output.getvalue()
    except Exception as e:
        logger.error(f"Error generating Excel report: {e}", exc_info=True)
        # Fallback: create an Excel file with an error message
        output_error = BytesIO()
        error_df = pd.DataFrame([{"Error": f"Could not generate Excel report: {str(e)}"}])
        with pd.ExcelWriter(output_error, engine="xlsxwriter") as writer:
            error_df.to_excel(writer, index=False, sheet_name="ErrorReport")
        return output_error.getvalue()

# --- END OF YOUR EXISTING USER FUNCTIONS ---


# ---- WORKER FUNCTION for Background Processing ----
# This function MUST be defined at the top level of this script file.
def worker_process_file(file_content_bytes, original_filename, file_mime_type, selected_classifications_list):
    """
    Processes a single file in a background process.
    All its dependencies (other functions it calls) must also be top-level or importable.
    """
    process_id = os.getpid()
    logger.info(f"PID {process_id}: Worker started for '{original_filename}'")
    
    file_buffer = BytesIO(file_content_bytes)
    ext = os.path.splitext(original_filename)[1].lower()
    
    raw_text = ""
    final_extracted_details_dict = {}
    masked_details_display_str = "No specific details extracted or targeted."
    # These are the final values for the report row
    report_doc_name = "Unknown" 
    report_doc_type = "Unclassified"
    report_sensitivity = "Low"
    report_is_pii = False

    try:
        # --- 1. Text Extraction ---
        if ext in [".png", ".jpg", ".jpeg"]:
            pil_img = Image.open(file_buffer)
            raw_text = extract_text_easyocr(pil_img) # Calls load_easyocr_model
        elif ext == ".pdf" or file_mime_type == "application/pdf":
            raw_text = extract_data_from_pdf(file_buffer)
        elif ext in [".txt", ".py", ".java", ".c", ".cpp", ".js", ".json", ".xml"]: # Assuming these are plain text
            raw_text = extract_text_plain(file_buffer)
        elif ext == ".pptx" or file_mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            raw_text = extract_text_from_pptx(file_buffer)
        elif ext == ".docx" or file_mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            raw_text = extract_text_from_docx(file_buffer)
        elif ext in [".xlsx", ".csv"] or file_mime_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "text/csv"]:
            df = process_excel_csv(file_buffer) # This might raise ValueError
            columns = df.columns.tolist()
            # For spreadsheets, details are usually columns; specific PII detection would be more complex (classify_csv_columns)
            report_doc_name = "Spreadsheet/CSV"
            report_doc_type = "Structured Data"
            report_sensitivity = "Potentially High (Manual Review Recommended)" # Default for spreadsheets
            masked_details_display_str = f"Columns: {', '.join(columns)}. (Content requires targeted PII scan)"
            report_is_pii = True # Assume PII might be present in any spreadsheet

            report_row = {
                "document_name": original_filename, "document": report_doc_name,
                "details": masked_details_display_str, "sensitivity": report_sensitivity,
                "document_type": report_doc_type, "is_pii": report_is_pii
            }
            logger.info(f"PID {process_id}: Worker successfully processed Spreadsheet/CSV '{original_filename}'")
            return {"status": "success", "filename": original_filename, "report_row": report_row, "error_message": None}
        else:
            logger.warning(f"PID {process_id}: Unsupported file extension '{ext}' for '{original_filename}'")
            raw_text = f"Error: Unsupported file type '{ext}'." # Set raw_text to an error message

        # --- 2. Preprocessing (if text was extracted or is an error message) ---
        if "Error:" in raw_text: # If extraction failed, raw_text contains the error
            cleaned_text = raw_text # Keep error message as is for now
        else:
            cleaned_text = preprocess_text(raw_text)

        text_upper = cleaned_text.upper() if cleaned_text else ""
        
        # --- 3. Document Identification ---
        identified_doc_primary_type = identify_document_type(text_upper)
        logger.info(f"PID {process_id}: Identified '{original_filename}' (type: {ext}) as '{identified_doc_primary_type}' based on content.")
        
        # --- 4. Specific Details Extraction (if identified type is in user's selection) ---
        if identified_doc_primary_type != "Unknown" and identified_doc_primary_type in selected_classifications_list:
            extraction_functions_map = {
                "Aadhaar": extract_aadhaar_details, "PAN": extract_pan_details,
                "Voter ID": extract_voter_details, "Passport": extract_passport_details, # Add your other extract functions
                "Driving License": extract_driving_license_details,
                "Medical Document": extract_medical_details,
                "Financial Document": extract_financial_details
            }
            if identified_doc_primary_type in extraction_functions_map:
                final_extracted_details_dict = extraction_functions_map[identified_doc_primary_type](cleaned_text)
                # Add the identified type to the details dict for consistency
                final_extracted_details_dict["Document_Title_Identified"] = identified_doc_primary_type
            
            if final_extracted_details_dict and "Error" not in final_extracted_details_dict:
                 masked_details_display_str = str(mask_data(final_extracted_details_dict))
            elif final_extracted_details_dict and "Error" in final_extracted_details_dict: # Extraction function returned an error
                 masked_details_display_str = f"Error in {identified_doc_primary_type} extraction: {final_extracted_details_dict['Error']}"
            elif identified_doc_primary_type != "Unknown": # Identified but no specific fields parsed by chosen extractor
                 masked_details_display_str = f"Generic '{identified_doc_primary_type}' document; no specific fields targeted by extractor or extracted."
        elif "Error:" in cleaned_text: # If extraction itself failed
            masked_details_display_str = cleaned_text # Show the extraction error

        # --- 5. Document Classification (based on extracted details and text) ---
        # The classify_document function should ideally use the final_extracted_details_dict
        report_doc_name, report_doc_type, report_sensitivity = classify_document(final_extracted_details_dict, text_upper)
        report_is_pii = report_doc_name not in ["Unknown", "Other"] and report_sensitivity in ["Medium", "High"] # Adjust PII logic

        report_row = {
            "document_name": original_filename, "document": report_doc_name,
            "details": masked_details_display_str, "sensitivity": report_sensitivity,
            "document_type": report_doc_type, "is_pii": report_is_pii
        }
        logger.info(f"PID {process_id}: Worker successfully processed '{original_filename}'")
        return {"status": "success", "filename": original_filename, "report_row": report_row, "error_message": None}

    except Exception as e:
        logger.error(f"PID {process_id}: Unhandled worker error for '{original_filename}': {str(e)}", exc_info=True)
        # Create an error report row
        error_report_row = {
            "document_name": original_filename, "document": "Processing Error",
            "details": f"Critical error in worker: {str(e)}", "sensitivity": "Unknown",
            "document_type": "Error", "is_pii": False 
        }
        return {"status": "error", "filename": original_filename, "report_row": error_report_row, "error_message": str(e)}


# ---- Helper function to manage ProcessPoolExecutor ----
def get_app_executor():
    """
    Safely gets or creates the ProcessPoolExecutor from/to st.session_state.
    Ensures only one executor instance per Streamlit session.
    """
    # Use a more specific key for this app's executor in session_state
    executor_session_key = "doc_scanner_process_pool_executor"
    if executor_session_key not in st.session_state:
        cpu_cores = os.cpu_count()
        if cpu_cores is None or cpu_cores <= 1: num_workers = 1
        elif cpu_cores <= 4: num_workers = 2 # For 2-4 cores, use 2 workers
        else: num_workers = min(cpu_cores -1, 4)  # Max 4 workers, or num_cores-1 for >4 cores
                                                # Adjust this based on OCR memory usage and typical file count
        
        logger.info(f"Main App (PID {os.getpid()}): Initializing ProcessPoolExecutor with {num_workers} workers.")
        st.session_state[executor_session_key] = ProcessPoolExecutor(max_workers=num_workers)
        
        # Optional: Add a cleanup for the executor when the session ends (if Streamlit version supports it)
        # This is advanced and depends on Streamlit's evolving session management.
        # For now, OS typically cleans up child processes when main Streamlit process exits.
        # if hasattr(st, 'on_event'): # Fictional example
        #     def _shutdown_executor():
        #         if executor_session_key in st.session_state:
        #             logger.info("Attempting to shutdown executor on session end.")
        #             st.session_state[executor_session_key].shutdown(wait=True) # wait=False for quicker exit
        #             del st.session_state[executor_session_key]
        #     st.on_event("session_end", _shutdown_executor)

    return st.session_state[executor_session_key]

# ---- Main Streamlit Application ----
def main():
    # Initialize or get the executor at the start of every script run.
    # `get_app_executor` handles creating it once per session via `st.session_state`.
    app_executor = get_app_executor()

    st.set_page_config(page_title="Document Scanner & Classifier", layout="wide")
    st.sidebar.title("🗂️ Document Intelligence")
    
    page_options = ["File Scanner", "DB Scanner (Placeholder)"]
    page_selection_icon = "📄" # Default icon
    if "File Scanner" in page_options[0]: page_selection_icon = "📄"
    elif "DB Scanner" in page_options[0]: page_selection_icon = "🗄️"
        
    page = st.sidebar.radio("Navigation", page_options, 
                            # format_func=lambda x: f"{'📄' if 'File' in x else '🗄️'} {x}" # If you want icons in radio
                           )


    # Initialize session state for tracking tasks and aggregated results
    if "file_processing_status_v2" not in st.session_state: # Use a distinct key
        st.session_state.file_processing_status_v2 = {}
    if "final_report_rows_collected_v2" not in st.session_state:
        st.session_state.final_report_rows_collected_v2 = []
    
    if page == "File Scanner":
        st.header("🔍 File Scanner & Document Classifier")
        st.markdown("Upload your documents (images, PDFs, Office files, text) for automated text extraction, PII identification, and classification.")

        with st.form(key="file_upload_form"):
            uploaded_files = st.file_uploader(
                "Choose files to scan",
                type=["png", "jpg", "jpeg", "pdf", "txt", "xlsx", "docx", "csv", "pptx"],
                accept_multiple_files=True,
                help="You can upload multiple files of supported types."
            )
            selected_classifications = st.multiselect(
                "Target specific document types for detailed extraction:",
                options=list(classification_options.keys()), # Your existing options
                default=list(classification_options.keys()),
                help="Select which document types you want the system to try and extract specific fields from."
            )
            scan_button_pressed = st.form_submit_button(label="🚀 Start Scanning Files", use_container_width=True)

        if scan_button_pressed:
            if uploaded_files:
                files_newly_queued_count = 0
                for up_file_obj in uploaded_files:
                    # Use Streamlit's built-in file_id for unique tracking if available (Streamlit >= 1.11.0)
                    # Fallback to name + size for older versions or if file_id is None
                    file_unique_id = getattr(up_file_obj, 'file_id', None)
                    if file_unique_id is None:
                        file_unique_id = f"{up_file_obj.name}_{up_file_obj.size}"
                    
                    # Process if file is new or previously errored (allowing retry)
                    task_state = st.session_state.file_processing_status_v2.get(file_unique_id)
                    if not task_state or task_state["status"] == "error":
                        logger.info(f"Main App: Queuing '{up_file_obj.name}' (ID: {file_unique_id}) for processing.")
                        try:
                            file_bytes_content = up_file_obj.getvalue() # Read bytes in main thread
                            
                            future_obj = app_executor.submit(
                                worker_process_file, # Target worker function
                                file_bytes_content,
                                up_file_obj.name,
                                up_file_obj.type, # Pass MIME type as a hint
                                selected_classifications # Pass user's selection list
                            )
                            st.session_state.file_processing_status_v2[file_unique_id] = {
                                "filename": up_file_obj.name, "status": "queued", 
                                "future": future_obj, "result_data": None,
                                "submitted_ts": datetime.datetime.now().isoformat() # For sorting/tracking
                            }
                            files_newly_queued_count += 1
                        except Exception as submission_err: # Catch errors during file reading or task submission
                            logger.error(f"Main App: Error submitting '{up_file_obj.name}' to executor: {submission_err}", exc_info=True)
                            st.session_state.file_processing_status_v2[file_unique_id] = {
                                "filename": up_file_obj.name, "status": "error", 
                                "future": None, "result_data": {"error_message": f"File submission failed: {str(submission_err)}"},
                                "submitted_ts": datetime.datetime.now().isoformat()
                            }
                
                if files_newly_queued_count > 0:
                    st.toast(f"{files_newly_queued_count} file(s) added to the scanning queue.", icon="⏳")
                    st.rerun() # Rerun to update UI immediately with "queued" statuses
                elif not uploaded_files: # Should be caught by outer if, but defensive
                    st.warning("Please upload at least one file to start scanning.")
                else: # All uploaded files were likely already processed or are in queue
                    st.info("All selected files are already in the queue or have been processed. You can clear previous results to rescan.")
            else: # No files uploaded but scan button pressed
                st.warning("No files uploaded. Please choose files to scan.")

        # --- Displaying Progress and Results for All Tasks ---
        active_tasks_exist = False # Flag to trigger polling if any task is still "processing"
        if st.session_state.file_processing_status_v2:
            st.write("---") # Visual separator
            st.subheader("📊 File Processing Dashboard")

            # Sort tasks by submission time for a consistent display order
            sorted_task_file_ids = sorted(
                st.session_state.file_processing_status_v2.keys(),
                key=lambda fid: st.session_state.file_processing_status_v2[fid]["submitted_ts"]
            )

            # Use columns for a more organized layout if many files
            # num_cols = 3 # Or 2 or 4
            # display_cols = st.columns(num_cols)
            # current_col_idx = 0

            for file_id in sorted_task_file_ids:
                task_info = st.session_state.file_processing_status_v2[file_id]
                
                # Update status from "queued" to "processing" once future is set (mainly visual)
                if task_info["status"] == "queued" and task_info.get("future"):
                    task_info["status"] = "processing"

                # Determine if expander should be open by default (e.g., for active or errored tasks)
                expander_open_by_default = task_info["status"] in ["queued", "processing", "error"]
                
                # current_display_container = display_cols[current_col_idx % num_cols]
                # with current_display_container: # For column layout

                with st.expander(f"📁 {task_info['filename']}  |  Status: {task_info['status'].upper()}", expanded=expander_open_by_default):
                    st.caption(f"Submitted: {task_info['submitted_ts']}")
                    
                    if task_info["status"] == "processing":
                        active_tasks_exist = True # At least one task is running
                        st.progress(50) # Indeterminate progress bar for processing
                        # st.spinner("Scanning in progress...") # Spinner can be redundant with progress

                        # Check if the future (background task) has completed
                        if task_info["future"] and task_info["future"].done():
                            needs_ui_refresh = False
                            try:
                                result_dict = task_info["future"].result(timeout=0.2) # Short timeout, should be done
                                task_info["result_data"] = result_dict
                                
                                if result_dict.get("status") == "success":
                                    task_info["status"] = "completed"
                                    # Add to overall report if not a duplicate (check by filename)
                                    if result_dict.get("report_row") and not any(
                                        r["document_name"] == result_dict["filename"] for r in st.session_state.final_report_rows_collected_v2
                                    ):
                                        st.session_state.final_report_rows_collected_v2.append(result_dict["report_row"])
                                    st.success(f"'{task_info['filename']}' processed successfully.")
                                else: # Worker explicitly reported an error
                                    task_info["status"] = "error"
                                    # If worker returns a report_row even on error, add it
                                    if result_dict.get("report_row") and not any(
                                        r["document_name"] == result_dict["filename"] for r in st.session_state.final_report_rows_collected_v2):
                                         st.session_state.final_report_rows_collected_v2.append(result_dict["report_row"])
                                    st.warning(f"Processing of '{task_info['filename']}' encountered an issue. Details below.")
                                needs_ui_refresh = True
                            except FutureTimeoutError:
                                logger.debug(f"Future for {task_info['filename']} timed out on result(), still processing.")
                                # This is normal, means task is not yet done, polling will catch it.
                            except Exception as result_retrieval_err: # Error during .result() call itself
                                logger.error(f"Main App: Critical error getting result for '{task_info['filename']}': {result_retrieval_err}", exc_info=True)
                                task_info["status"] = "error"
                                task_info["result_data"] = {"error_message": f"System error retrieving result: {str(result_retrieval_err)}"}
                                # Add a generic error row to the main report for tracking
                                generic_error_row = {
                                    "document_name": task_info['filename'], "document": "System Error",
                                    "details": f"Failed to retrieve processing result: {str(result_retrieval_err)}", 
                                    "sensitivity": "Unknown", "document_type": "Error", "is_pii": False
                                }
                                if not any(r["document_name"] == task_info['filename'] for r in st.session_state.final_report_rows_collected_v2):
                                    st.session_state.final_report_rows_collected_v2.append(generic_error_row)
                                needs_ui_refresh = True
                            
                            if needs_ui_refresh:
                                st.rerun() # Rerun to update UI immediately after a task changes state

                    # Display final status details
                    if task_info["status"] == "completed":
                        st.success("✅ Processing complete.")
                        # st.json(task_info["result_data"].get("report_row", {})) # Optional: show individual report row details here
                    elif task_info["status"] == "error":
                        st.error("❌ An error occurred during processing:")
                        st.text_area("Error Details:", 
                                     value=str(task_info.get("result_data", {}).get("error_message", "No specific error message available.")), 
                                     height=100, disabled=True)
                    elif task_info["status"] == "queued":
                        st.info("⏳ Waiting in queue to be processed...")
                # current_col_idx +=1


            # --- Display Overall Consolidated Report Table ---
            if st.session_state.final_report_rows_collected_v2:
                st.write("---") # Visual separator
                st.subheader("📋 Consolidated Scan Report")
                
                # Allow user to filter report (example: by sensitivity)
                # sensitivity_filter = st.multiselect("Filter by Sensitivity:", 
                #                                     options=df_report['sensitivity'].unique() if not df_report.empty else [], 
                #                                     key="report_sensitivity_filter")
                # if sensitivity_filter:
                #     df_to_display = df_report[df_report['sensitivity'].isin(sensitivity_filter)]
                # else:
                #     df_to_display = df_report
                
                df_report = pd.DataFrame(st.session_state.final_report_rows_collected_v2)
                st.dataframe(df_report, use_container_width=True)

                # Summary Metrics
                pii_files_count = sum(1 for row in st.session_state.final_report_rows_collected_v2 if row.get("is_pii", False))
                total_files_in_report = len(st.session_state.final_report_rows_collected_v2)
                
                metric_cols = st.columns(2)
                metric_cols[0].metric("Total Files in Report", total_files_in_report)
                metric_cols[1].metric("Files Classified with PII", pii_files_count, delta_color="inverse") # Or "off"

                # Download Button
                try:
                    excel_data_bytes = save_to_excel(st.session_state.final_report_rows_collected_v2)
                    st.download_button(
                        label="📥 Download Full Report (Excel)",
                        data=excel_data_bytes,
                        file_name="Comprehensive_Document_Scan_Report.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        key="download_consolidated_report_button",
                        use_container_width=True
                    )
                except Exception as excel_gen_err:
                    st.error(f"Could not generate Excel report for download: {excel_gen_err}")
            
            # If there are tasks still actively processing, schedule a UI refresh (polling)
            if active_tasks_exist:
                # Make polling interval slightly longer to reduce busy-waiting feel
                time.sleep(5) # e.g., 5 seconds
                logger.debug(f"Main App: Active tasks detected ({sum(1 for t in st.session_state.file_processing_status_v2.values() if t['status']=='processing')}). Triggering UI refresh for polling.")
                st.rerun()
        
        # Button to clear all scanned data and the current report
        st.write("---")
        if st.button("🧹 Clear All Scanned Data & Reset Report", key="clear_all_scan_data_button", use_container_width=True):
            # Note: This doesn't stop already running background processes, just clears UI state.
            # True cancellation of ProcessPoolExecutor tasks is complex and often not fully reliable.
            st.session_state.file_processing_status_v2 = {}
            st.session_state.final_report_rows_collected_v2 = []
            # Attempt to clear the file uploader if it has a key (might not always work perfectly across Streamlit versions)
            if "file_uploader_widget" in st.session_state:
                st.session_state.file_uploader_widget = [] 
            st.toast("All scanned file data and the current report have been cleared.", icon="♻️")
            st.rerun()


    elif page == "DB Scanner (Placeholder)": # DB Scanner Page
        st.header("🗄️ Database Scanner")
        st.warning("This feature is currently under development. Please check back later.", icon="⚠️")
        
        # Placeholder for DB connection UI
        # st.subheader("Connect to Database")
        # dialect = st.selectbox("Database Dialect", ["postgresql", "mysql", "sqlite", "sqlserver", "oracle"], key="db_dialect")
        # username = st.text_input("Username", key="db_user")
        # password = st.text_input("Password", type="password", key="db_pass")
        # host = st.text_input("Host", value="localhost", key="db_host")
        # database = st.text_input("Database Name", key="db_name")
        # port = st.text_input("Port", value="5432", key="db_port") # Default for PostgreSQL
        # if st.button("Test Connection & Scan", key="db_connect_scan_button"):
        #     st.info(f"Attempting to connect to {dialect} at {host}:{port}...")
        #     # ... (Your DB connection and scanning logic would go here) ...
        #     st.success("Connection successful! (Scanning logic not yet implemented).")


if __name__ == "__main__":
    # This guard is CRITICAL for ProcessPoolExecutor to work correctly,
    # especially on Windows. It prevents child processes from re-executing
    # the main script when they are spawned.
    main()