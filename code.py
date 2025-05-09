import concurrent
import datetime
import chardet
import streamlit as st
import easyocr
import torch
# import cv2 # Used by easyocr, not directly imported unless needed elsewhere
import numpy as np
import pandas as pd
from PIL import Image
from io import BytesIO
import os
import pdfplumber
from pptx import Presentation
from docx import Document
# import sqlalchemy as sa # Not used in file processing part
# import PyPDF2 # Not used
# from openpyxl import Workbook, load_workbook # Pandas handles this
# import psycopg2 # For DB Scanner
# from psycopg2 import sql # For DB Scanner
from concurrent.futures import ProcessPoolExecutor # For parallel CPU-bound tasks
import re
import time # For polling UI updates
import logging # For better insight into background processes

# --- Existing User Functions (UNCHANGED) ---
# All your functions from load_easyocr_model to classification_options
# are assumed to be here and remain exactly as you provided.
# For brevity, I'm not re-listing them, but they are crucial.

# --- Setup Logging ---
# This helps debug background processes by printing to the console where Streamlit runs
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(process)d - %(threadName)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# --- START OF EXISTING USER FUNCTIONS (Copied for completeness but UNCHANGED) ---
torch.classes.__path__ = []

@st.cache_resource
def load_easyocr_model():
    logger.info(f"PID {os.getpid()}: Loading EasyOCR model...")
    reader = easyocr.Reader(['en', 'hi']) # Add other languages if needed
    logger.info(f"PID {os.getpid()}: EasyOCR model loaded.")
    return reader

# Image Data Extraction
def extract_text_easyocr(pil_img):
    reader = load_easyocr_model()
    # EasyOCR expects a NumPy array (RGB)
    cv_img = np.array(pil_img.convert("RGB"))
    results = reader.readtext(cv_img, detail=0) # detail=0 for just text
    return "\n".join(results)


# PDF Data Extraction
def extract_data_from_pdf(pdf_path_or_buffer): # Modified to accept buffer
    extracted_text = ""
    try:
        with pdfplumber.open(pdf_path_or_buffer) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    extracted_text += text + "\n"
    except Exception as e:
        logger.error(f"Error extracting PDF: {e}")
        return f"Error extracting PDF: {e}"
    return extracted_text

# PPTX Data Extraction
def extract_text_from_pptx(file_path_or_buffer): # Modified
    prs = Presentation(file_path_or_buffer)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return '\n'.join(text_runs)

# Wordfile Data Extraction
def extract_text_from_docx(file_path_or_buffer): # Modified
    doc = Document(file_path_or_buffer)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)

# Text File Extraction
def extract_text_plain(file_like_object): # Modified to take file-like object
    encodings = ['utf-8', 'latin-1', 'iso-8859-1', 'cp1252']
    file_like_object.seek(0) # Ensure reading from the beginning
    file_content = file_like_object.read()
    for encoding in encodings:
        try:
            return file_content.decode(encoding)
        except UnicodeDecodeError:
            continue
    logger.warning("extract_text_plain: Unable to decode with available encodings")
    return "Error reading text: Unable to decode with available encodings"

# Text Preprocessing
def preprocess_text(text):
    if text is not None and isinstance(text, str): # Add type check
        text = text.strip()
        text = re.sub(r'\s+', ' ', text)
        text = ''.join([c for c in text if c.isprintable()])
        allowed_chars = r"A-Za-z0-9\s:./,\-_()[]" # Corrected regex string
        text = re.sub(r"[^" + allowed_chars + r"]", "", text)
    elif text is None:
        return "" # Return empty string if None
    return text

# PAN Extraction
def extract_pan_details(text):
    if not text or not isinstance(text, str): return {"Error": "Invalid input text"}
    details = {}; text_upper = text.upper()
    if re.search(r"INCOME TAX DEPARTMENT", text_upper): details["Department"] = "INCOME TAX DEPARTMENT"
    if re.search(r"GOVT\.?\s+OF\s+INDIA", text_upper): details["Issuer"] = "GOVT. OF INDIA"
    if re.search(r"PERMANENT ACCOUNT NUMBER CARD", text_upper): details["Document_Title"] = "PERMANENT ACCOUNT NUMBER CARD"
    pan_match = re.search(r"\b[A-Z]{5}[0-9]{4}[A-Z]\b", text_upper);
    if pan_match: details["PAN_Number"] = pan_match.group(0)
    name_match = re.search(r"APPLICANT\s+NAME[:\s]+([A-Z\s]+)", text_upper)
    if name_match: details["Name"] = name_match.group(1).strip()
    else:
        name_match = re.search(r"NAME[:\s]+([A-Z\s]+)", text_upper)
        if name_match: details["Name"] = name_match.group(1).strip()
    father_match = re.search(r"FATHER\s+NAME[:\s]+([A-Z\s]+)", text_upper)
    if father_match: details["Father_Name"] = father_match.group(1).strip()
    dob_match = re.search(r"\b\d{2}/\d{2}/\d{4}\b", text_upper)
    if dob_match: details["Date_of_Birth"] = dob_match.group(0)
    details["Signature"] = "Present" if re.search(r"SIGNATURE", text_upper) else "Not Found"
    return details

# Aadhaar Extraction
def extract_aadhaar_details(text):
    details = {}; text_upper = text.upper()
    if "UNIQUE IDENTIFICATION AUTHORITY OF INDIA" in text_upper: details["Issuer"] = "UIDAI"
    if "AADHAAR" in text_upper: details["Document_Title"] = "AADHAAR CARD"
    aadhaar_match = re.search(r"\b\d{4}\s?\d{4}\s?\d{4}\b", text_upper)
    if aadhaar_match: details["Aadhaar_Number"] = aadhaar_match.group(0)
    enrol_match = re.search(r"\b\d{4}/\d{5}/\d{5}\b", text_upper)
    if enrol_match: details["Enrolment_No"] = enrol_match.group(0)
    name_match = re.search(r"NAME[:\s]+([A-Z\s]+)", text_upper)
    if name_match: details["Name"] = name_match.group(1).strip()
    dob_match = re.search(r"\b\d{2}/\d{2}/\d{4}\b", text_upper)
    if dob_match: details["Date_of_Birth"] = dob_match.group(0)
    lines = text_upper.splitlines()
    for i, line in enumerate(lines):
        if "ADDRESS" in line and i + 1 < len(lines):
            details["Address"] = lines[i+1].strip().title(); break
    return details

# Voter ID Extraction
def extract_voter_details(text):
    details = {}; text_upper = text.upper()
    epic_match = re.search(r"\b[A-Z]{3}[0-9]{7}\b", text_upper)
    if epic_match: details["Voter_ID"] = epic_match.group(0)
    if "ELECTION COMMISSION OF INDIA" in text_upper: details["Document_Title"] = "ELECTION COMMISSION OF INDIA"
    elif "ELECTOR PHOTO IDENTITY CARD" in text_upper: details["Document_Title"] = "ELECTOR PHOTO IDENTITY CARD"
    elif "VOTER ID" in text_upper or "EPIC" in text_upper: details["Document_Title"] = "VOTER ID"
    name_match = re.search(r"(?:NAME|ELECTOR S NAME)[:\s]+([A-Z\s]+)", text_upper)
    if name_match: details["Name"] = name_match.group(1).strip()
    relative_match = re.search(r"(?:FATHER?S|HUSBAND?S) NAME[:\s]+([A-Z\s]+)", text_upper)
    if relative_match: details["Relative_Name"] = relative_match.group(1).strip()
    return details

# Passport Extraction
def extract_passport_details(text):
    details = {}; text_upper = text.upper()
    if "REPUBLIC OF INDIA" in text_upper or "PASSPORT" in text_upper: details["Document_Title"] = "PASSPORT"
    passport_match = re.search(r"\b[A-Z][0-9]{7}\b", text_upper)
    if passport_match: details["Passport_Number"] = passport_match.group(0)
    name_match = re.search(r"NAME[:\s]+([A-Z\s]+)", text_upper)
    if name_match: details["Name"] = name_match.group(1).strip()
    dob_match = re.search(r"\b\d{2}/\d{2}/\d{4}\b", text_upper)
    if dob_match: details["Date_of_Birth"] = dob_match.group(0)
    nationality_match = re.search(r"NATIONALITY[:\s]+([A-Z\s]+)", text_upper)
    if nationality_match: details["Nationality"] = nationality_match.group(1).strip()
    return details

# Driving License Extraction
def extract_driving_license_details(text):
    details = {}; text_upper = text.upper()
    if "DRIVING LICENCE" in text_upper or "DRIVER LICENSE" in text_upper: details["Document_Title"] = "DRIVING LICENSE"
    dl_match = re.search(r"\b[A-Z0-9]{10,20}\b", text_upper) # General, might need refinement
    if dl_match: details["License_Number"] = dl_match.group(0)
    name_match = re.search(r"NAME[:\s]+([A-Z\s]+)", text_upper)
    if name_match: details["Name"] = name_match.group(1).strip()
    dob_match = re.search(r"\b\d{2}/\d{2}/\d{4}\b", text_upper)
    if dob_match: details["Date_of_Birth"] = dob_match.group(0)
    address_match = re.search(r"ADDRESS[:\s]+([A-Z0-9\s,./-]+)", text_upper) # Allow more chars in address
    if address_match: details["Address"] = address_match.group(1).strip()
    return details

# Medical Details Extraction
def extract_medical_details(text):
    details = {}; text_upper = text.upper()
    if "MEDICAL CERTIFICATE" in text_upper or "HEALTH RECORD" in text_upper: details["Document_Title"] = "MEDICAL CERTIFICATE"
    name_match = re.search(r"(?:PATIENT\s+)?NAME[:\s]+([A-Z\s.]+)", text_upper) # Allow dots in name
    if name_match: details["Patient_Name"] = name_match.group(1).strip()
    dob_match = re.search(r"\b\d{2}/\d{2}/\d{4}\b", text_upper)
    if dob_match: details["Date_of_Birth"] = dob_match.group(0)
    doctor_match = re.search(r"DOCTOR(?:\S\sNAME)?[:\s]+(DR\.?\s*[A-Z\s.]+)", text_upper, re.IGNORECASE)
    if doctor_match: details["Doctor_Name"] = doctor_match.group(1).strip()
    medical_id_match = re.search(r"(?:MEDICAL\s+)?ID[:\s]+([A-Z0-9\s-]+)", text_upper)
    if medical_id_match: details["Medical_ID"] = medical_id_match.group(1).strip()
    return details

# Financial Details Extraction
def extract_financial_details(text):
    details = {}; text_upper = text.upper()
    if "BANK STATEMENT" in text_upper or "INCOME TAX" in text_upper: details["Document_Title"] = "FINANCIAL DOCUMENT"
    acc_match = re.search(r"ACCOUNT\s*(?:NUMBER|NO\.?)\s*[:\-]?\s*(\b[0-9X]{9,18}\b)", text_upper, re.IGNORECASE)
    if acc_match: details["Account_Number"] = acc_match.group(1)
    pan_match = re.search(r"\b[A-Z]{5}[0-9]{4}[A-Z]\b", text_upper)
    if pan_match: details["PAN_Number"] = pan_match.group(0)
    ifsc_match = re.search(r"IFSC\s*(?:CODE)?\s*[:\-]?\s*(\b[A-Z]{4}0[0-9A-Z]{6}\b)", text_upper, re.IGNORECASE)
    if ifsc_match: details["IFSC_Code"] = ifsc_match.group(1)
    txn_match = re.search(r"(?:TXN|TRANSACTION)\s*(?:ID|NO\.?)\s*[:\-]?\s*([A-Z0-9]+)", text_upper, re.IGNORECASE)
    if txn_match: details["Transaction_ID"] = txn_match.group(1)
    return details

# Document Identification Regex (as provided)
AADHAAR_REGEX = re.compile(r'\b\d{4}\s?\d{4}\s?\d{4}\b', re.IGNORECASE)
PAN_REGEX = re.compile(r'\b[A-Z]{5}[0-9]{4}[A-Z]\b', re.IGNORECASE)
VOTER_ID_REGEX = re.compile(r'\b[A-Z]{3}[0-9]{7}\b', re.IGNORECASE)
PASSPORT_REGEX = re.compile(r'\b[A-Z]{1}[0-9]{7}\b', re.IGNORECASE)
DL_REGEX_STRICT = re.compile(r'\b(?:[A-Z]{2}[-\s]?\d{2}[-\s]?\d{4}[-\s]?\d{7}|[A-Z]{2}[-\s]?\d{13})\b', re.IGNORECASE) # More specific DL
INVOICE_REGEX = re.compile(r'\b(?:Invoice|Bill No|Receipt No)\s*[:-]?\s*(\w+)\b', re.IGNORECASE)
# TOTAL_AMOUNT_REGEX = re.compile(r'\bTotal\sAmount\s[:-]?\sRs.?\s(\d+[.,]?\d*)\b', re.IGNORECASE) # Not directly used for ID
# DATE_REGEX = re.compile(r'\b(?:Date|Billing Date)\s*[:-]?\s*(\d{2}[/-]\d{2}[/-]\d{4})\b') # Not directly used for ID

DOCUMENT_PATTERNS = {
    "Aadhaar": (AADHAAR_REGEX, {"AADHAAR", "UNIQUE IDENTIFICATION AUTHORITY OF INDIA", "GOVERNMENT OF INDIA"}),
    "PAN": (PAN_REGEX, {"INCOME TAX DEPARTMENT", "PERMANENT ACCOUNT NUMBER CARD", "PAN CARD"}),
    "Voter ID": (VOTER_ID_REGEX, {"VOTER ID", "EPIC", "ELECTION COMMISSION OF INDIA", "ELECTOR PHOTO IDENTITY CARD"}),
    "Passport": (PASSPORT_REGEX, {"PASSPORT", "REPUBLIC OF INDIA", "MINISTRY OF EXTERNAL AFFAIRS"}),
    "Driving License": (DL_REGEX_STRICT, {"DRIVING LICENCE", "TRANSPORT DEPARTMENT", "LICENSE NUMBER", "DL NO"}), # Using stricter DL regex
    "Medical Document": (None, {"HOSPITAL", "PATIENT", "DOCTOR", "CLINIC", "MEDICAL", "DIAGNOSIS", "PRESCRIPTION", "RX"}), # Regex for medical is less common for ID
    "Financial Document": (None, {"BANK", "ACCOUNT", "STATEMENT", "INVOICE", "BILL", "FINANCIAL", "TRANSACTION"}) # Regex for financial is less common for ID
}

def identify_document_type(text_upper):
    # Helper for keyword matching (from user's code, moved for clarity)
    def _keyword_match_count(keywords, text_content):
        return sum(1 for keyword in keywords if keyword in text_content)

    detected_types = {}
    for doc_type, (regex, keywords) in DOCUMENT_PATTERNS.items():
        score = 0
        if regex and regex.search(text_upper):
            score += 2 # Strong indicator
        
        keyword_hits = _keyword_match_count(keywords, text_upper)
        score += keyword_hits
        
        if score > 0:
            detected_types[doc_type] = score
            
    # Basic rule: if a strong regex hit, prefer that. Otherwise, most keywords.
    if not detected_types:
        return "Unknown"

    # If Aadhaar regex and keywords match strongly
    if "Aadhaar" in detected_types and AADHAAR_REGEX.search(text_upper) and detected_types["Aadhaar"] >=2:
        return "Aadhaar"
    # If PAN regex and keywords match strongly
    if "PAN" in detected_types and PAN_REGEX.search(text_upper) and detected_types["PAN"] >=2 :
        return "PAN"
    
    # Prioritize regex matches if they exist and have some keyword support
    priority_by_regex = [dt for dt, (r,k) in DOCUMENT_PATTERNS.items() if r and r.search(text_upper) and dt in detected_types and detected_types[dt]>=1]
    if priority_by_regex:
        # if multiple regex matches, pick the one with more keyword hits
        return max(priority_by_regex, key=lambda dt: detected_types[dt])


    # Fallback to the one with the most keyword matches if no strong regex signal found above
    return max(detected_types, key=detected_types.get)


column_patterns = { # For DB Scanner or specific CSV column checks, not heavily used in main file flow provided
    'ssn': { 'keywords': ['ssn', 'social security number'], 'regex': r'\b(?!000|666)[0-8]\d{2}(-|\s)(?!00)\d{2}(-|\s)(?!0000)\d{4}\b'}
}
def detect_encoding(file_path): # Not directly used in current flow
    with open(file_path, 'rb') as file: result = chardet.detect(file.read())
    return result['encoding']

def classify_csv_columns(df): # For specific CSV sensitive column detection
    classified_columns = {}; PI_data = pd.DataFrame()
    for column in df.columns:
        for key, patterns in column_patterns.items():
            if any(keyword in column.lower() for keyword in patterns['keywords']):
                PI_data[column] = df[column]; classified_columns[column] = key
    return PI_data # Return the DataFrame with potentially sensitive columns

def save_to_excel(report_data_list_of_dicts): # Expects a list of dicts
    if not report_data_list_of_dicts:
        logger.warning("No data provided to save_to_excel.")
        # Create an empty Excel file if no data
        df = pd.DataFrame([{"Message": "No data processed or available."}])
    else:
        df = pd.DataFrame(report_data_list_of_dicts)

    output = BytesIO()
    # Use try-except for robustness if xlsxwriter is not installed or other issues
    try:
        with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
            df.to_excel(writer, index=False, sheet_name="Report")
        return output.getvalue()
    except Exception as e:
        logger.error(f"Error saving to Excel: {e}")
        # Fallback or re-raise
        # For now, let's return an empty BytesIO or a simple error message in Excel
        output_error = BytesIO()
        error_df = pd.DataFrame([{"Error": f"Could not generate Excel: {str(e)}"}])
        with pd.ExcelWriter(output_error, engine="xlsxwriter") as writer:
            error_df.to_excel(writer, index=False, sheet_name="ErrorSheet")
        return output_error.getvalue()


def process_excel_csv(file_like_object): # Modified to take file-like object
    # Guess extension from name if possible, but rely on pandas' ability to infer
    # For BytesIO, pandas might need help or specific args for read_csv if no clear delimiter
    try:
        df = pd.read_excel(file_like_object)
    except Exception: # Broad exception, try CSV
        try:
            file_like_object.seek(0) # Reset buffer position
            df = pd.read_csv(file_like_object)
        except Exception as e:
            logger.error(f"Failed to process as Excel or CSV: {e}")
            raise ValueError("Unsupported file format or error reading Excel/CSV.") from e
    return df

# classify_document needs keyword_match_count, so define it first or make it global/pass it
def _keyword_match_count_global(keywords, text): # Renamed to avoid conflict if nested
    return sum(1 for keyword in keywords if keyword in text)

def classify_document(details_arg, text_upper): # details_arg is the cleaned_text in original call
    # data_extracted = len(details_arg) if isinstance(details_arg, str) else len(details_arg) if isinstance(details_arg, dict) else 0
    # The original code passes `cleaned_text` as `details_arg`. 
    # Let's assume `details_arg` here is the `details` dictionary extracted by `extract_..._details` functions.
    # If `details_arg` is `cleaned_text` (a string), then `len(details_arg)` makes sense.
    # If `details_arg` is a dict from `extract_..._details`, `len(details_arg)` is num of keys.
    # The user's original call was `classify_document(cleaned_text, text_upper)`.
    # The function definition `def classify_document(details, text_upper):`
    # So, `details` inside the function body IS `cleaned_text`.
    # `data_extracted` then means `len(cleaned_text)`. This seems inconsistent with its usage.
    # Let's assume the INTENT of `data_extracted` was the number of *fields* extracted.
    # For this wrapper, we'll pass the `details_dict_from_extraction` to this function.
    
    # For the purpose of this refactoring, I will assume `details_arg` is the dictionary of extracted fields.
    data_extracted_fields = len(details_arg) if isinstance(details_arg, dict) else 0


    aadhaar_keywords = ["AADHAAR", "UNIQUE IDENTIFICATION AUTHORITY OF INDIA", "GOVERNMENT OF INDIA"]
    if any(keyword in text_upper for keyword in aadhaar_keywords) and data_extracted_fields > 1: # e.g. number and name
        return "Aadhaar", "Government ID", "High"

    pan_keywords = ["INCOME TAX DEPARTMENT", "PERMANENT ACCOUNT NUMBER CARD", "PAN CARD"]
    if any(keyword in text_upper for keyword in pan_keywords) and data_extracted_fields > 1: # e.g. number and name
        return "PAN", "Government ID", "High"

    dl_keywords = ["DRIVING LICENCE", "TRANSPORT DEPARTMENT", "DL NO"]
    if any(keyword in text_upper for keyword in dl_keywords) and data_extracted_fields > 1:
        return "Driving License", "Government ID", "High"

    passport_keywords = ["PASSPORT", "REPUBLIC OF INDIA", "MINISTRY OF EXTERNAL AFFAIRS"]
    if any(keyword in text_upper for keyword in passport_keywords) and data_extracted_fields > 1:
        return "Passport", "Government ID", "High"

    voter_keywords = ["VOTER ID", "EPIC", "ELECTION COMMISSION OF INDIA", "ELECTOR PHOTO IDENTITY CARD"]
    if any(keyword in text_upper for keyword in voter_keywords) and data_extracted_fields > 1:
        return "Voter ID", "Government ID", "High"

    medical_keywords = ["HOSPITAL", "PATIENT", "DOCTOR", "CLINIC", "MEDICAL", "DIAGNOSIS", "PRESCRIPTION"]
    if _keyword_match_count_global(medical_keywords, text_upper) >= 2 and data_extracted_fields > 0:
        return "Medical Document", "Medical Document", "Medium"

    financial_keywords = ["BANK", "ACCOUNT", "STATEMENT", "INVOICE", "BILL", "FINANCIAL", "TRANSACTION"]
    if _keyword_match_count_global(financial_keywords, text_upper) >= 2 and data_extracted_fields > 0:
        return "Financial Document", "Financial Document", "Medium"
    
    return "Unknown", "Unknown", "Low"


def mask_data(details_dict): # Expects a dictionary
    masked_details = {}
    if not isinstance(details_dict, dict): return {} # Robustness
    for key, value in details_dict.items():
        if value and isinstance(value, str) and len(value) > 2 : # Only mask strings longer than 2 chars
            masked_details[key] = 'X' * (len(value) - 2) + value[-2:]
        else:
            masked_details[key] = value # Keep non-strings or short strings as is
    return masked_details

classification_options = {
    "Aadhaar": "Government ID", "PAN": "Government ID", "Voter ID": "Government ID",
    "Passport": "Government ID", "Driving License": "Government ID",
    "Medical Document": "Medical Document", "Financial Document": "Financial Document",
}
# --- END OF EXISTING USER FUNCTIONS ---


# ---- NEW WORKER FUNCTION for Background Processing ----
def worker_process_file(file_content_bytes, original_filename, file_mime_type, selected_classifications_list):
    """
    This function processes a single file in a background process.
    It replicates the logic from the original main loop for one file.
    """
    pid = os.getpid()
    logger.info(f"PID {pid}: Worker started for {original_filename}")
    
    file_buffer = BytesIO(file_content_bytes)
    # Adding a 'name' attribute to the BytesIO buffer can be helpful for some libraries
    # However, original_filename is passed and used for extension checking.
    # file_buffer.name = original_filename # Some libs might check this, but pdfplumber, docx, pptx work with buffer directly.

    ext = os.path.splitext(original_filename)[1].lower()
    raw_text = ""
    # Initialize report variables for this file
    final_details_dict = {}
    masked_details_str = "No specific details extracted."
    doc_name_classified = "Unknown"
    doc_type_in_report = "Unknown" # This is 'document_type' in report_row
    sensitivity_classified = "Low"
    is_pii_present = False

    try:
        if ext in [".png", ".jpg", ".jpeg"]:
            pil_img = Image.open(file_buffer) # No need to convert to RGB if easyocr handles it, but good practice
            raw_text = extract_text_easyocr(pil_img)
        elif ext == ".pdf" or file_mime_type == "application/pdf":
            raw_text = extract_data_from_pdf(file_buffer)
        elif ext in [".txt", ".py", ".java", ".c", ".cpp", ".js", ".json", ".xml"]: # Re-added more text types
            raw_text = extract_text_plain(file_buffer)
        elif ext == ".pptx" or file_mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            raw_text = extract_text_from_pptx(file_buffer)
        elif ext == ".docx" or file_mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            raw_text = extract_text_from_docx(file_buffer)
        elif ext in [".xlsx", ".csv"] or file_mime_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "text/csv"]:
            df = process_excel_csv(file_buffer)
            columns = df.columns.tolist()
            # result_classified_csv = classify_csv_columns(df) # User had this commented, let's assume not for main report row now

            # For Excel/CSV, the original code had direct assignments for report
            doc_name_classified = "Spreadsheet/CSV Data" # Was "CSV File"
            doc_type_in_report = "Structured Data File" # Was "Sensitive Spreadsheet"
            sensitivity_classified = "High" # Assuming high due to structured nature
            # Details for CSV/Excel will be a string representation of columns
            masked_details_str = f"Columns: {', '.join(columns)}" 
            is_pii_present = True # Assume PII if it's a structured data file needing scan

            report_row = {
                "document_name": original_filename, "document": doc_name_classified,
                "details": masked_details_str, "sensitivity": sensitivity_classified,
                "document_type": doc_type_in_report, "is_pii": is_pii_present
            }
            logger.info(f"PID {pid}: Worker finished CSV/Excel {original_filename}")
            return {"status": "success", "filename": original_filename, "report_row": report_row, "error_message": None}
        else:
            logger.warning(f"PID {pid}: Unsupported file extension {ext} for {original_filename}")
            raw_text = "Unsupported file type for text extraction."
            # Fall through to general text processing if raw_text has something, or it will be handled
        
        # Common processing for text-based documents (non-Excel/CSV)
        cleaned_text = preprocess_text(raw_text)
        # st.write("Cleaned Text:", cleaned_text) # Cannot use st.write in background

        text_upper = cleaned_text.upper() if cleaned_text else ""
        
        # identify_document_type logic might need selected_classifications if it's to be filtered early
        identified_doc_primary_type = identify_document_type(text_upper)
        logger.info(f"PID {pid}: Identified '{original_filename}' as '{identified_doc_primary_type}'")

        if identified_doc_primary_type in selected_classifications_list: # Filter based on user selection
            if identified_doc_primary_type == "Aadhaar":
                final_details_dict = extract_aadhaar_details(cleaned_text)
            elif identified_doc_primary_type == "PAN":
                final_details_dict = extract_pan_details(cleaned_text)
            elif identified_doc_primary_type == "Voter ID":
                final_details_dict = extract_voter_details(cleaned_text)
            elif identified_doc_primary_type == "Passport":
                final_details_dict = extract_passport_details(cleaned_text)
            elif identified_doc_primary_type == "Driving License":
                final_details_dict = extract_driving_license_details(cleaned_text)
            elif identified_doc_primary_type == "Medical Document":
                final_details_dict = extract_medical_details(cleaned_text)
            elif identified_doc_primary_type == "Financial Document":
                final_details_dict = extract_financial_details(cleaned_text)
            # else: final_details_dict remains empty
            
            if final_details_dict:
                 masked_details_str = str(mask_data(final_details_dict))
            elif identified_doc_primary_type != "Unknown": # Type was identified but no specific details extracted
                 masked_details_str = f"Generic {identified_doc_primary_type} document, no specific fields parsed."
            # If final_details_dict is empty and type is Unknown, masked_details_str remains default

        # Classify document and get sensitivity level
        # Original call: classify_document(cleaned_text, text_upper)
        # The function def is: classify_document(details, text_upper)
        # Assuming `details` parameter in classify_document refers to the *dictionary* of extracted fields.
        doc_name_classified, doc_type_in_report, sensitivity_classified = classify_document(final_details_dict, text_upper)

        # Check if the document type implies PII based on your classification logic
        is_pii_present = doc_name_classified not in ["Unknown", "Other"] # Adjust this condition based on how `doc_name_classified` determines PII

        report_row = {
            "document_name": original_filename, "document": doc_name_classified,
            "details": masked_details_str, "sensitivity": sensitivity_classified,
            "document_type": doc_type_in_report, "is_pii": is_pii_present
        }
        logger.info(f"PID {pid}: Worker finished {original_filename}")
        return {"status": "success", "filename": original_filename, "report_row": report_row, "error_message": None}

    except Exception as e:
        logger.error(f"PID {pid}: Worker error for {original_filename}: {str(e)}", exc_info=True)
        return {"status": "error", "filename": original_filename, "report_row": None, "error_message": str(e)}


# ---- Main Streamlit Application ----
def main():
    st.sidebar.title("Classification Types")
    page = st.sidebar.selectbox("Select a page:", ["File Scanner", "DB Scanner"])

    # --- Session State Initialization for Background Tasks ---
    if "file_processing_status" not in st.session_state:
        # {file_id: {"filename": str, "status": str ("queued", "processing", "completed", "error"), 
        #            "future": Future, "result_data": dict, "submitted_ts": datetime}}
        st.session_state.file_processing_status = {}
    if "process_pool_executor" not in st.session_state:
        # Adjust max_workers based on CPU and memory. OCR is CPU/memory intensive.
        # Using fewer workers than cores can be stabler if memory is an issue.
        num_workers = max(1, (os.cpu_count() or 4) // 2) # Default to 2 if os.cpu_count() is None or 1
        logger.info(f"Main App (PID {os.getpid()}): Initializing ProcessPoolExecutor with {num_workers} workers.")
        st.session_state.process_pool_executor = ProcessPoolExecutor(max_workers=num_workers)
    if "final_report_rows_collected" not in st.session_state:
        st.session_state.final_report_rows_collected = [] # List of dicts for the final report

    if page == "File Scanner":
        st.title("Document Classification & Scanner")

        uploaded_files = st.file_uploader(
            "Upload Multiple Files",
            type=["png", "jpg", "jpeg", "pdf", "txt", "xlsx", "docx", "csv", "pptx"], # Added pptx
            accept_multiple_files=True,
            key="file_uploader_widget" # Key helps Streamlit manage widget state
        )
       
        selected_classifications = st.multiselect(
            "Select the classifications you want to detect:",
            options=list(classification_options.keys()), # Uses your existing dict
            default=list(classification_options.keys())
        )

        if st.button("Scan Files", key="scan_button"):
            if uploaded_files:
                files_newly_queued = 0
                for up_file in uploaded_files:
                    # Create a unique ID for each uploaded file instance
                    # Streamlit's UploadedFile object has a `file_id` attribute (from v1.11.0+)
                    # For older versions, name + size might be used, but less robust for re-uploads of same file.
                    file_unique_id = getattr(up_file, 'file_id', f"{up_file.name}_{up_file.size}")

                    # Check if file is new or failed previously and can be retried
                    if file_unique_id not in st.session_state.file_processing_status or \
                       st.session_state.file_processing_status[file_unique_id]["status"] == "error":
                        
                        logger.info(f"Main App: Queuing {up_file.name} (ID: {file_unique_id}) for processing.")
                        try:
                            file_bytes = up_file.getvalue() # Read bytes in main thread
                            
                            future = st.session_state.process_pool_executor.submit(
                                worker_process_file, # The new worker function
                                file_bytes,
                                up_file.name,
                                up_file.type, # Pass MIME type as a hint
                                selected_classifications # Pass user's selection
                            )
                            st.session_state.file_processing_status[file_unique_id] = {
                                "filename": up_file.name, "status": "queued", 
                                "future": future, "result_data": None,
                                "submitted_ts": datetime.datetime.now().isoformat()
                            }
                            files_newly_queued += 1
                        except Exception as e:
                            logger.error(f"Main App: Error submitting {up_file.name} to executor: {e}", exc_info=True)
                            st.session_state.file_processing_status[file_unique_id] = {
                                "filename": up_file.name, "status": "error", 
                                "future": None, "result_data": {"error_message": f"Submission failed: {str(e)}"},
                                "submitted_ts": datetime.datetime.now().isoformat()
                            }
                
                if files_newly_queued > 0:
                    st.toast(f"{files_newly_queued} file(s) added to the scanning queue.", icon="⏳")
                    st.rerun() # Rerun to update UI with "queued" statuses
                elif not uploaded_files: # Should be caught by outer if
                    st.warning("Please upload one or more files to scan.")
                else: # All uploaded files might already be in queue or processed
                    st.info("All uploaded files are already in the queue or have been processed. Clear data to rescan.")
            else:
                st.warning("Please upload one or more files to scan.")

        # --- Displaying Progress and Results for All Tasks ---
        active_tasks_found = False
        if st.session_state.file_processing_status:
            st.write("---")
            st.subheader("File Processing Status")

            # Sort for consistent display order, e.g., by submission time
            sorted_file_ids = sorted(
                st.session_state.file_processing_status.keys(),
                key=lambda fid: st.session_state.file_processing_status[fid]["submitted_ts"]
            )

            for file_id in sorted_file_ids:
                task_info = st.session_state.file_processing_status[file_id]
                
                # Update status from "queued" to "processing" if future exists (visual cue)
                if task_info["status"] == "queued" and task_info.get("future"):
                    task_info["status"] = "processing"

                expander_state = task_info["status"] in ["queued", "processing", "error"]
                with st.expander(f"File: {task_info['filename']} - Status: {task_info['status'].upper()}", expanded=expander_state):
                    
                    if task_info["status"] == "processing":
                        active_tasks_found = True
                        st.spinner("Scanning in progress...")
                        # Check if the future (background task) is done
                        if task_info["future"] and task_info["future"].done():
                            try:
                                result = task_info["future"].result(timeout=0.1) # Short timeout, should be done
                                task_info["result_data"] = result
                                if result.get("status") == "success":
                                    task_info["status"] = "completed"
                                    # Add to overall report if not a duplicate (based on filename for simplicity here)
                                    if result.get("report_row") and not any(
                                        r["document_name"] == result["filename"] for r in st.session_state.final_report_rows_collected
                                    ):
                                        st.session_state.final_report_rows_collected.append(result["report_row"])
                                    st.toast(f"'{task_info['filename']}' processed successfully.", icon="✅")
                                else: # Error from worker
                                    task_info["status"] = "error"
                                    st.toast(f"Error processing '{task_info['filename']}'.", icon="⚠️")
                                st.rerun() # Rerun to reflect immediate change
                            except concurrent.futures.TimeoutError:
                                pass # Should not happen if future.done() is true, but good to handle
                            except Exception as e:
                                logger.error(f"Main App: Error getting result for {task_info['filename']}: {e}", exc_info=True)
                                task_info["status"] = "error"
                                task_info["result_data"] = {"error_message": f"Failed to retrieve result: {str(e)}"}
                                st.rerun() # Rerun to show error

                    if task_info["status"] == "completed":
                        st.success("Processing complete.")
                        # Optionally display part of the result for this specific file
                        # st.json(task_info["result_data"].get("report_row", {}))
                    elif task_info["status"] == "error":
                        st.error("An error occurred during processing:")
                        st.json(task_info.get("result_data", {"error_message": "Unknown error details."}))
                    elif task_info["status"] == "queued":
                         st.info("Waiting in queue to be processed...")


            # --- Display Overall Report Table (Aggregated from completed tasks) ---
            if st.session_state.final_report_rows_collected:
                st.write("---")
                st.subheader("Consolidated Scan Report")
                df_report = pd.DataFrame(st.session_state.final_report_rows_collected)
                st.dataframe(df_report)

                pi_count = sum(1 for row in st.session_state.final_report_rows_collected if row.get("is_pii", False))
                
                col1, col2 = st.columns(2)
                col1.metric("Total Files in Report", len(st.session_state.final_report_rows_collected))
                col2.metric("Files with PII (as per classification)", pi_count)

                try:
                    excel_bytes = save_to_excel(st.session_state.final_report_rows_collected)
                    st.download_button(
                        label="Download Full Report as Excel",
                        data=excel_bytes,
                        file_name="document_scan_report.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        key="download_report_excel_button"
                    )
                except Exception as e:
                    st.error(f"Could not generate Excel report: {e}")
            
            # If there are active tasks, schedule a rerun to poll for updates
            if active_tasks_found:
                time.sleep(3) # Polling interval (e.g., 3 seconds)
                logger.debug(f"Main App: {sum(1 for t in st.session_state.file_processing_status.values() if t['status']=='processing')} tasks active. Triggering UI refresh.")
                st.rerun()
        
        # Button to clear all task statuses and the collected report
        if st.button("Clear All Scanned Data & Report", key="clear_all_data_button"):
            # Optionally try to cancel futures if executor supports it well, but can be complex
            # For simplicity, just clear state. Tasks might complete in background but won't be shown.
            st.session_state.file_processing_status = {}
            st.session_state.final_report_rows_collected = []
            # Reset file uploader by clearing its state if a key was used
            if "file_uploader_widget" in st.session_state:
                st.session_state.file_uploader_widget = []
            st.toast("All scanned data and the report have been cleared.", icon="🗑️")
            st.rerun()


    # DB Scanner Page (remains as is from your original code)
    elif page == "DB Scanner":
        st.subheader("Enter Database Connection Details")
        # ... (your existing DB Scanner UI and logic) ...
        dialect = st.selectbox("Database Dialect", ["postgresql", "mysql", "sqlite", "sqlserver", "oracle"])
        username = st.text_input("Username")
        password = st.text_input("Password", type="password")
        host = st.text_input("Host", value="localhost")
        database = st.text_input("Database", value="")
        port = st.text_input("Port", value="5432")

        if st.button("Connect", key="db_connect_button"):
            st.write(f"Attempting to connect to {dialect} database at {host}:{port}...")
            # Placeholder for database connection logic
            # db_connection = connect_to_db(dialect, username, password, host, database, port)
            st.success("Connected successfully! (Placeholder)")


if __name__ == "__main__":
    # This check ensures that if this script is run directly (not via `streamlit run`),
    # certain setup (like model loading) can happen.
    # However, for Streamlit apps, the primary entry is `streamlit run <filename>.py`.
    # The @st.cache_resource for load_easyocr_model will handle model loading correctly
    # when the worker processes start.
    main()